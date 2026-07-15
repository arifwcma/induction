import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from build_deck import (
    BG, PANEL, PANEL_LINE, TEAL, AMBER, WHITE, MUTED, WRONG, FOOTER_BG,
    FONT, MONO, SLIDE_W, SLIDE_H,
    solid, no_line, outline, txt, card, circle, base_slide,
)


HERE = os.path.dirname(os.path.abspath(__file__))
OUT_PATH = os.path.join(HERE, "chunking_deck.pptx")

EYEBROW = "PHASE 3  ·  SEARCH FIRST"
TITLE = "Chunking — what we gain, what we lose"


def bullet_column(slide, x, w, heading, heading_color, items):
    top = 2.25
    height = 4.35
    card(slide, x, top, w, height, fill=PANEL, line_color=heading_color)

    badge = circle(slide, x + 0.48, top + 0.46, 0.46, fill=heading_color)
    sign = "+" if heading_color == TEAL else "\u2212"
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sign
    run.font.name = FONT
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = BG

    txt(slide, x + 0.95, top + 0.2, w - 1.15, 0.55, [(heading, 16, heading_color, True, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    item_y = top + 1.02
    for title_text, detail_text in items:
        txt(slide, x + 0.4, item_y, w - 0.75, 0.32, [(title_text, 12.5, WHITE, True, False, None)])
        txt(slide, x + 0.4, item_y + 0.32, w - 0.75, 0.5, [(detail_text, 10, MUTED, False, False, 1.1)])
        item_y += 0.8


def build_slide(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES)

    txt(slide, 0.9, 1.82, 11.6, 0.4,
        [("Sending a few chunks instead of the whole KB is a trade-off — cheaper and faster, but not free.",
          13.5, MUTED, False, True, None)])

    pros = [
        ("Small, fast, cheap prompt",
         "a few paragraphs instead of 2 million words — answers in seconds, at a fraction of the cost."),
        ("Sharper focus",
         "less irrelevant text means less to distract the model — often a more accurate answer."),
        ("Scales to a huge KB",
         "hundreds of documents become searchable; only the winners ever reach the model."),
        ("Fits the context limit",
         "a model can only read so much at once — chunking keeps us safely under that ceiling."),
    ]
    cons = [
        ("Missing the holistic view",
         "the answer may live across several chunks; picking a few can sever the bigger picture."),
        ("A weak retriever bottlenecks a strong model",
         "a cheap/local search may drop a chunk that GPT-level model could have used well — the model never sees what it never receives."),
        ("Wrong chunk \u2192 confident wrong answer",
         "if search picks the wrong passage, the model answers fluently from it (our library 24/7 bug)."),
        ("Boundaries split ideas",
         "a clumsy cut can slice one rule in half, leaving each piece meaningless on its own."),
    ]

    gap = 0.5
    col_w = (SLIDE_W - 0.9 - 0.9 - gap) / 2
    bullet_column(slide, 0.9, col_w, "Pros", TEAL, pros)
    bullet_column(slide, 0.9 + col_w + gap, col_w, "Cons", WRONG, cons)

    txt(slide, 0.9, 6.68, 11.6, 0.32,
        [("The craft is choosing chunks well — so the pros stay and the cons shrink.", 12, AMBER, False, True, None)],
        align=PP_ALIGN.CENTER)
    return slide


NOTES = (
    "Chunking is the trade we make in phase three, and it is worth being honest about both sides. "
    "On the plus side: the prompt becomes small, fast and cheap; the model is less distracted so answers "
    "can be sharper; it scales to hundreds of documents; and it keeps us under the model's context limit. "
    "But there is a cost. We can lose the holistic view - the real answer might be spread across several "
    "chunks, and by sending only a few we cut the bigger picture. The retriever also becomes a bottleneck: "
    "if we use a weak or local search to pick chunks, it may throw away a passage that a stronger model "
    "like GPT would have found useful - the model can only reason over what we actually hand it. If search "
    "picks the wrong chunk, the model still answers confidently from it - exactly our library 24/7 bug. And "
    "clumsy boundaries can split a single rule in half. So the whole skill is choosing chunks well: keep "
    "the speed and cost wins, shrink the losses."
)


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_slide(prs)
    prs.save(OUT_PATH)
    print("saved", OUT_PATH)


if __name__ == "__main__":
    main()
