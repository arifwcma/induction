import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


BG = RGBColor(0x0B, 0x1B, 0x33)
PANEL = RGBColor(0x14, 0x28, 0x47)
PANEL_LINE = RGBColor(0x25, 0x3B, 0x5E)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
AMBER = RGBColor(0xF5, 0xB3, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9F, 0xB3, 0xC8)
WRONG = RGBColor(0xE0, 0x6C, 0x6C)
FOOTER_BG = RGBColor(0x0F, 0x22, 0x40)

FONT = "Segoe UI"
MONO = "Consolas"
SLIDE_W = 13.333
SLIDE_H = 7.5

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(HERE, "lecture_resources", "logo-wimmera-white.png")
OUT_PATH = os.path.join(HERE, "deck.pptx")

SPEAKER_NAME = "Dr Mohammad Arifur Rahman"
SPEAKER_ROLE = "Analyst Programmer  ·  Wimmera CMA, Victoria"


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def outline(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def set_letterspacing(run, points):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(points * 100)))


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, size, color, bold, italic, spacing) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return box


def card(slide, x, y, w, h, fill=PANEL, line_color=PANEL_LINE, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    solid(shape, fill)
    outline(shape, line_color, 1.0)
    shape.adjustments[0] = radius
    return shape


def frame(slide, x, y, w, h, line_color, width_pt=1.5, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    outline(shape, line_color, width_pt)
    shape.adjustments[0] = radius
    return shape


def circle(slide, cx, cy, diameter, fill=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - diameter / 2), Inches(cy - diameter / 2),
                                    Inches(diameter), Inches(diameter))
    solid(shape, fill)
    no_line(shape)
    return shape


def centered_label(shape, text, size, color, bold=True, font=FONT):
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def arrow_right(slide, x, y, w, h=0.3, fill=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y - h / 2), Inches(w), Inches(h))
    solid(shape, fill)
    no_line(shape)
    return shape


def straight_line(slide, x1, y1, x2, y2, color, width_pt):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)
    return shape


def add_logo(slide, height, top=0.35, right_margin=0.45):
    if not os.path.exists(LOGO_PATH):
        return
    pic = slide.shapes.add_picture(LOGO_PATH, Inches(0), Inches(top), height=Inches(height))
    pic.left = Inches(SLIDE_W - right_margin) - pic.width


def add_footer(slide, prs):
    fh = 0.55
    fy = SLIDE_H - fh
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(fy), prs.slide_width, Inches(fh))
    solid(footer, FOOTER_BG)
    no_line(footer)
    faccent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(fy), prs.slide_width, Inches(0.04))
    solid(faccent, TEAL)
    no_line(faccent)
    txt(slide, 0.65, fy, 5.0, fh, [("ICT108  ·  MetMate case study", 10, MUTED, False, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)
    right_box = slide.shapes.add_textbox(Inches(5.8), Inches(fy), Inches(SLIDE_W - 5.8 - 0.45), Inches(fh))
    tf = right_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r1 = p.add_run()
    r1.text = SPEAKER_NAME
    r1.font.name = FONT
    r1.font.size = Pt(10.5)
    r1.font.bold = True
    r1.font.color.rgb = TEAL
    r2 = p.add_run()
    r2.text = "   ·   " + SPEAKER_ROLE
    r2.font.name = FONT
    r2.font.size = Pt(10)
    r2.font.color.rgb = MUTED
    return fy


def base_slide(prs, eyebrow_label, title, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    solid(bg, BG)
    no_line(bg)

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), prs.slide_height)
    solid(strip, TEAL)
    no_line(strip)

    eb = txt(slide, 0.65, 0.4, 9.5, 0.4, [(eyebrow_label, 12.5, TEAL, True, False, None)])
    set_letterspacing(eb.text_frame.paragraphs[0].runs[0], 2.2)

    txt(slide, 0.63, 0.74, 10.6, 0.9, [(title, 30, WHITE, True, False, 1.0)])

    ul = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(1.5), Inches(1.3), Inches(0.06))
    solid(ul, AMBER)
    no_line(ul)

    add_logo(slide, 0.5)
    add_footer(slide, prs)

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def failure_chip(slide, text, x=0.9, y=1.85, w=4.0):
    chip = card(slide, x, y, w, 0.52, fill=PANEL, line_color=WRONG)
    txt(slide, x + 0.18, y, w - 0.36, 0.52, [(text, 11.5, WRONG, True, True, None)], anchor=MSO_ANCHOR.MIDDLE)
    return chip


def term_chip(slide, x, y, w, h, term, definition, def_size=10.5):
    chip = card(slide, x, y, w, h, fill=PANEL, line_color=TEAL)
    box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y), Inches(w - 0.4), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.line_spacing = 1.12
    r1 = p.add_run()
    r1.text = term
    r1.font.name = FONT
    r1.font.size = Pt(12.5)
    r1.font.bold = True
    r1.font.color.rgb = TEAL
    r2 = p.add_run()
    r2.text = "  —  " + definition
    r2.font.name = FONT
    r2.font.size = Pt(def_size)
    r2.font.color.rgb = MUTED
    return chip


def flow_box(slide, x, y, w, h, label, line_color, size=11.5, text_color=WHITE, fill=PANEL):
    shape = card(slide, x, y, w, h, fill=fill, line_color=line_color)
    txt(slide, x + 0.08, y, w - 0.16, h, [(label, size, text_color, True, False, 1.08)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return shape


def flow_row(slide, left, right, y, h, items, gap=0.26, size=11.5, arrow_color=TEAL):
    n = len(items)
    w = (right - left - gap * (n - 1)) / n
    positions = []
    for i, (label, line_color) in enumerate(items):
        x = left + i * (w + gap)
        flow_box(slide, x, y, w, h, label, line_color, size=size)
        positions.append((x, w))
        if i < n - 1:
            arrow_right(slide, x + w + 0.02, y + h / 2, gap - 0.04, h=0.18, fill=arrow_color)
    return positions


def build_s01_title(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    solid(bg, BG)
    no_line(bg)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), prs.slide_height)
    solid(strip, TEAL)
    no_line(strip)

    eb = txt(slide, 0.65, 1.25, 11.0, 0.4,
             [("GUEST LECTURE   ·   ICT108   ·   SYDNEY MET", 13, TEAL, True, False, None)])
    set_letterspacing(eb.text_frame.paragraphs[0].runs[0], 2.4)

    txt(slide, 0.62, 1.7, 12.0, 1.3, [("Talk to Your Documents", 46, WHITE, True, False, 1.0)])

    ul = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(2.72), Inches(1.9), Inches(0.08))
    solid(ul, AMBER)
    no_line(ul)

    txt(slide, 0.65, 3.0, 11.8, 0.6,
        [("How a chatbot that answers from documents really works — built from zero, phase by phase.",
          17, MUTED, False, False, 1.1)])

    txt(slide, 0.65, 3.85, 11.5, 0.45, [("Case study: MetMate", 15, TEAL, True, False, None)])
    txt(slide, 0.65, 4.28, 11.5, 0.45,
        [("An invented chatbot answering Sydney Met student questions — enrolment, library, exams.",
          13, MUTED, False, True, None)])

    promise = card(slide, 0.65, 5.0, 10.2, 0.95, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 0.92, 5.0, 9.7, 0.95,
        [("Concepts over terminology: every technical term arrives only when the story needs it,", 12, MUTED, False, True, 1.25),
         ("defined in plain English. Nothing to memorise up front — everything is look-up-able later.", 12, MUTED, False, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    txt(slide, 0.65, 6.25, 12.2, 0.4,
        [("Phase 1  Just ask ChatGPT   →   Phase 2  Our own bot   →   Phase 3  Search first   →   "
          "Phase 4  Search by meaning   →   Phase 5  Toward the real world", 10.5, TEAL, False, False, None)])

    add_logo(slide, 1.0, top=0.45)
    add_footer(slide, prs)

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def build_s02_upload(prs, notes):
    slide = base_slide(prs, "PHASE 1  ·  JUST ASK CHATGPT", "Upload a PDF, ask a question", notes)

    txt(slide, 0.9, 1.85, 11.6, 0.4,
        [("The simplest possible start — no code at all.", 13.5, MUTED, False, True, None)])

    chat = card(slide, 0.9, 2.4, 3.5, 3.5, fill=PANEL, line_color=TEAL)
    bar = card(slide, 0.9, 2.4, 3.5, 0.45, fill=FOOTER_BG, line_color=TEAL)
    txt(slide, 1.1, 2.4, 3.1, 0.45, [("ChatGPT", 12.5, TEAL, True, False, None)], anchor=MSO_ANCHOR.MIDDLE)
    pill = card(slide, 1.15, 3.05, 2.1, 0.42, fill=FOOTER_BG, line_color=PANEL_LINE, radius=0.5)
    txt(slide, 1.35, 3.05, 1.8, 0.42, [("handbook.pdf", 10.5, WHITE, False, False, None)], anchor=MSO_ANCHOR.MIDDLE)
    q = card(slide, 1.15, 3.65, 3.0, 0.8, fill=PANEL, line_color=TEAL, radius=0.25)
    txt(slide, 1.35, 3.65, 2.6, 0.8, [("What are the library opening hours?", 11, WHITE, False, False, 1.15)],
        anchor=MSO_ANCHOR.MIDDLE)
    a = card(slide, 1.15, 4.6, 3.0, 0.8, fill=FOOTER_BG, line_color=PANEL_LINE, radius=0.25)
    txt(slide, 1.35, 4.6, 2.6, 0.8, [("8am–10pm on term days. ✓", 11, MUTED, False, False, 1.15)],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, 0.9, 6.0, 3.5, 0.4, [("one person, one chat — it works", 11, TEAL, False, True, None)],
        align=PP_ALIGN.CENTER)

    txt(slide, 4.55, 3.6, 1.0, 0.35, [("under the", 10, MUTED, False, True, None)], align=PP_ALIGN.CENTER)
    txt(slide, 4.55, 3.85, 1.0, 0.35, [("hood", 10, MUTED, False, True, None)], align=PP_ALIGN.CENTER)
    arrow_right(slide, 4.55, 4.45, 0.85, fill=TEAL)

    big = card(slide, 5.55, 2.4, 4.0, 3.5, fill=PANEL, line_color=AMBER)
    txt(slide, 5.8, 2.55, 3.5, 0.4, [("what the model actually receives", 11, AMBER, True, False, None)])
    txt(slide, 5.8, 3.05, 3.5, 2.4, [
        ("<entire text of handbook.pdf>", 11.5, MUTED, False, True, 1.25),
        ("… 200 pages of it …", 11.5, MUTED, False, True, 1.25),
        ("+", 13, WHITE, True, False, 1.25),
        ("\u201cWhat are the library", 11.5, WHITE, False, False, 1.2),
        ("opening hours?\u201d", 11.5, WHITE, False, False, None),
    ], font=MONO)
    txt(slide, 5.8, 5.45, 3.5, 0.4, [("ONE long piece of text", 12, TEAL, True, False, None)])

    arrow_right(slide, 9.7, 4.15, 0.7, fill=TEAL)

    model = card(slide, 10.55, 3.3, 2.15, 1.7, fill=PANEL, line_color=TEAL)
    txt(slide, 10.7, 3.55, 1.85, 0.4, [("the model", 13, TEAL, True, False, None)], align=PP_ALIGN.CENTER)
    txt(slide, 10.7, 4.0, 1.85, 0.9, [("reads it all, writes the answer", 11, WHITE, False, False, 1.2)],
        align=PP_ALIGN.CENTER)

    note = card(slide, 4.55, 6.15, 8.15, 0.6, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 4.8, 6.15, 7.7, 0.6,
        [("Today: TEXT only. Images, audio, complex layouts — a story for another day.", 11, MUTED, False, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)
    return slide


def build_s03_doesnt_scale(prs, notes):
    slide = base_slide(prs, "PHASE 1  ·  IT DOESN'T SCALE", "Why this doesn't scale", notes)

    txt(slide, 0.9, 1.85, 11.6, 0.4,
        [("It works for one person. Now give it to 30,000 students.", 13.5, MUTED, False, True, None)])

    problems = [
        ("It's ChatGPT's product",
         "We can't put ChatGPT itself on the university website. Building our own ChatGPT-like product from scratch is a huge job."),
        ("A paid seat per student",
         "The alternative: every student buys their own subscription — thousands of seats, just to ask about the handbook."),
        ("Many documents, always changing",
         "One answer may need handbook + library rules + exam procedures. They change all year. Re-uploading the current pile into every chat? Nobody will."),
    ]
    left = 0.9
    gap = 0.25
    w = (SLIDE_W - 0.9 - left - 2 * gap) / 3
    for i, (head, body) in enumerate(problems):
        x = left + i * (w + gap)
        card(slide, x, 2.4, w, 2.45, fill=PANEL, line_color=PANEL_LINE)
        num = circle(slide, x + 0.42, 2.82, 0.44, fill=AMBER)
        centered_label(num, str(i + 1), 14, BG)
        txt(slide, x + 0.75, 2.6, w - 0.95, 0.75, [(head, 13.5, WHITE, True, False, 1.1)])
        txt(slide, x + 0.25, 3.45, w - 0.5, 1.3, [(body, 11, MUTED, False, False, 1.25)])

    term_chip(slide, 0.9, 5.15, 11.55, 0.75,
              "Knowledge Base (KB)",
              "the pile of official documents the bot must answer from. Our first term — "
              "every teal box like this one is a new term entering the story.")

    txt(slide, 0.9, 6.2, 11.6, 0.5,
        [("What we want: ONE bot, OURS, on the university site — always reading the CURRENT documents. So we build.",
          13.5, AMBER, False, True, None)])
    return slide


def build_s04_own_bot(prs, notes):
    slide = base_slide(prs, "PHASE 2  ·  OUR OWN BOT", "Our own bot — same trick, our pipes", notes)

    flow_row(slide, 0.9, 12.43, 2.3, 1.05, [
        ("Student's browser\nchat page", PANEL_LINE),
        ("Our server", PANEL_LINE),
        ("OpenAI API\nthe model behind ChatGPT", TEAL),
    ], size=12.5)

    txt(slide, 0.9, 3.5, 3.8, 0.35, [("ordinary web dev", 10, MUTED, False, True, None)], align=PP_ALIGN.CENTER)
    txt(slide, 4.96, 3.5, 3.8, 0.35, [("ordinary backend", 10, MUTED, False, True, None)], align=PP_ALIGN.CENTER)
    txt(slide, 9.02, 3.5, 3.5, 0.35, [("their computers, our text", 10, MUTED, False, True, None)], align=PP_ALIGN.CENTER)

    payload = card(slide, 3.3, 4.05, 6.7, 1.15, fill=PANEL, line_color=AMBER)
    txt(slide, 3.55, 4.2, 6.2, 0.4, [("what our server sends — the phase-1 trick, through pipes:", 11, AMBER, True, False, None)])
    txt(slide, 3.55, 4.6, 6.2, 0.55, [
        ("every KB document  +  the student's question", 12, WHITE, False, False, 1.2),
        ("\u2192  one big string", 12, TEAL, True, False, None),
    ], font=MONO)

    term_chip(slide, 0.9, 5.5, 7.1, 0.85,
              "API",
              "a doorway for programs: our code sends text over the internet, the model's reply "
              "comes back. ChatGPT is the product — OpenAI is the company and the API.")

    note = card(slide, 8.25, 5.5, 4.2, 0.85, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 8.5, 5.5, 3.7, 0.85,
        [("Local option: run a free model on your own machine (e.g. Ollama). Same story either way.",
          10.5, MUTED, False, True, 1.2)], anchor=MSO_ANCHOR.MIDDLE)

    txt(slide, 0.9, 6.55, 11.6, 0.4,
        [("For the rest of the talk: we use the OpenAI API.", 12, TEAL, False, True, None)])
    return slide


def build_s05_prompt(prs, notes):
    slide = base_slide(prs, "PHASE 2  ·  OUR OWN BOT", "What the model actually receives", notes)

    box = card(slide, 0.9, 2.0, 6.9, 4.25, fill=FOOTER_BG, line_color=TEAL)
    txt(slide, 1.35, 2.0, 6.0, 4.25, [
        ("Context:", 17, TEAL, True, False, 1.45),
        ("    <every document in the KB>", 15, MUTED, False, True, 1.45),
        ("", 10, MUTED, False, False, 1.0),
        ("Question:", 17, TEAL, True, False, 1.45),
        ("    <the student's question>", 15, MUTED, False, True, 1.45),
        ("", 10, MUTED, False, False, 1.0),
        ("Answer using only the", 17, WHITE, False, False, 1.45),
        ("context above.", 17, WHITE, False, False, None),
    ], font=MONO, anchor=MSO_ANCHOR.MIDDLE)

    term_chip(slide, 8.1, 2.0, 4.35, 1.0,
              "Prompt", "the full text we send to the model. This is literally everything it sees.")
    term_chip(slide, 8.1, 3.2, 4.35, 1.0,
              "Context", "the reference text we paste into the prompt, for the model to answer from.")

    note = card(slide, 8.1, 4.4, 4.35, 1.85, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 8.35, 4.6, 3.85, 1.5, [
        ("No memory of Sydney Met.", 11.5, WHITE, False, False, 1.35),
        ("No hidden database.", 11.5, WHITE, False, False, 1.35),
        ("If the right fact is not in the prompt, the model can only guess.", 11.5, MUTED, False, False, 1.25),
    ])

    txt(slide, 0.9, 6.5, 11.6, 0.45,
        [("The whole craft of a document chatbot: deciding WHAT goes into this text.", 13.5, AMBER, False, True, None)])
    return slide


def build_s06_ai_vs_not(prs, notes):
    slide = base_slide(prs, "PHASE 2  ·  AI VS NOT-AI", "How much of this is actually AI?", notes)

    y = 2.55
    h = 1.0
    flow_box(slide, 1.15, y, 2.2, h, "Student's browser", PANEL_LINE, size=11.5)
    arrow_right(slide, 3.42, y + h / 2, 0.34, h=0.18, fill=TEAL)
    flow_box(slide, 3.83, y, 2.2, h, "Our server", PANEL_LINE, size=11.5)
    arrow_right(slide, 6.1, y + h / 2, 0.34, h=0.18, fill=TEAL)
    flow_box(slide, 6.51, y, 2.2, h, "Glue the big string", PANEL_LINE, size=11.5)
    arrow_right(slide, 8.78, y + h / 2, 0.34, h=0.18, fill=TEAL)
    flow_box(slide, 9.65, y, 2.55, h, "The model", AMBER, size=12.5)

    frame(slide, 0.95, 2.25, 7.95, 1.6, TEAL, width_pt=1.75)
    txt(slide, 0.95, 3.95, 7.95, 0.4, [("ORDINARY CODE — deterministic", 12, TEAL, True, False, None)],
        align=PP_ALIGN.CENTER)
    frame(slide, 9.45, 2.25, 2.95, 1.6, AMBER, width_pt=1.75)
    txt(slide, 9.45, 3.95, 2.95, 0.4, [("AI — non-deterministic", 12, AMBER, True, False, None)],
        align=PP_ALIGN.CENTER)

    term_chip(slide, 0.9, 4.65, 5.65, 0.95,
              "Deterministic", "same input → same output, every time. Testable, predictable, debuggable.")
    term_chip(slide, 6.8, 4.65, 5.65, 0.95,
              "Non-deterministic", "same input → the output may vary. Ask twice, get two wordings.")

    txt(slide, 0.9, 5.9, 11.6, 0.5,
        [("An \u201cAI system\u201d is mostly NOT AI — one AI box, surrounded by ordinary software.", 15, AMBER, True, False, None)])
    txt(slide, 0.9, 6.45, 11.6, 0.4,
        [("Most of your work as a developer — and most of your value — lives in the deterministic part.", 12, MUTED, False, True, None)])
    return slide


def build_s07_too_slow(prs, notes):
    slide = base_slide(prs, "PHASE 2  ·  IT BREAKS", "One question, two million words", notes)
    failure_chip(slide, "\u2717 too slow to be usable", w=3.2)

    kb = card(slide, 0.9, 2.65, 3.4, 2.2, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.15, 2.85, 2.9, 0.4, [("the KB", 12, MUTED, True, False, None)])
    txt(slide, 1.15, 3.25, 2.9, 0.6, [("300 PDFs", 20, WHITE, True, False, None)])
    txt(slide, 1.15, 3.9, 2.9, 0.6, [("\u2248 2,000,000 words", 15, AMBER, True, False, None)])

    arrow_right(slide, 4.45, 3.75, 0.7, fill=TEAL)

    prompt = card(slide, 5.3, 2.65, 3.4, 2.2, fill=PANEL, line_color=AMBER)
    txt(slide, 5.55, 2.85, 2.9, 0.4, [("the prompt", 12, AMBER, True, False, None)])
    txt(slide, 5.55, 3.25, 2.9, 1.4, [
        ("ALL of it, for EVERY question", 13.5, WHITE, True, False, 1.3),
        ("even \u201cwhat are the library hours?\u201d", 11, MUTED, False, True, None),
    ])

    arrow_right(slide, 8.85, 3.75, 0.7, fill=TEAL)

    watch = card(slide, 9.7, 2.65, 2.95, 2.2, fill=PANEL, line_color=WRONG)
    txt(slide, 9.95, 2.85, 2.45, 0.4, [("the result", 12, WRONG, True, False, None)])
    txt(slide, 9.95, 3.25, 2.45, 0.6, [("~1 minute", 20, WRONG, True, False, None)])
    txt(slide, 9.95, 3.9, 2.45, 0.85, [("per answer — my first real bot worked exactly like this", 10.5, MUTED, False, True, 1.2)])

    txt(slide, 0.9, 5.1, 11.6, 0.4,
        [("Almost everything we send is irrelevant to the question being asked.", 12.5, MUTED, False, True, None)])

    magic = card(slide, 0.9, 5.65, 7.35, 1.05, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.15, 5.65, 6.85, 1.05,
        [("The model stays a MAGIC BOX today — how it works inside is a field of its own.", 11.5, WHITE, False, False, 1.25),
         ("Our job as chatbot developers: feed the box well — fast and accurate.", 11.5, MUTED, False, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    nxt = card(slide, 8.5, 5.65, 4.15, 1.05, fill=PANEL, line_color=TEAL)
    txt(slide, 8.75, 5.65, 3.65, 1.05,
        [("So: send LESS.", 13, TEAL, True, False, 1.25),
         ("But which part?  →  Phase 3", 12, WHITE, False, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)
    return slide


def build_s08_send_less(prs, notes):
    slide = base_slide(prs, "PHASE 3  ·  SEARCH FIRST", "Send less: pick the right pieces", notes)

    txt(slide, 0.9, 1.85, 11.6, 0.4,
        [("For any one question, only a few paragraphs of the whole KB matter. Send those.", 13.5, MUTED, False, True, None)])

    flow_row(slide, 0.9, 12.43, 2.45, 1.0, [
        ("KB\ndocuments", PANEL_LINE),
        ("split into\nchunks", TEAL),
        ("pick the top 5\nfor this question", TEAL),
        ("prompt:\n5 chunks + question", PANEL_LINE),
        ("model", PANEL_LINE),
        ("answer", PANEL_LINE),
    ], size=10.5)

    term_chip(slide, 0.9, 3.9, 7.1, 0.8,
              "Chunk", "a small passage cut from a document — a few paragraphs, ideally one topic.")

    shrink = card(slide, 8.25, 3.9, 4.2, 0.8, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 8.5, 3.9, 3.7, 0.8, [("200 pages  →  2 pages per question", 12.5, TEAL, True, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    skill = card(slide, 0.9, 5.05, 11.55, 1.45, fill=PANEL, line_color=AMBER)
    txt(slide, 1.2, 5.22, 11.0, 0.4, [("This choosing step is YOUR craft", 14, AMBER, True, False, None)])
    txt(slide, 1.2, 5.65, 11.0, 0.75,
        [("How the model answers from text — not our department, we take it as given. Choosing WHAT it reads — that is "
          "exactly where a chatbot developer shows their skill. Most of what made my real bot better happened right here.",
          11.5, WHITE, False, False, 1.25)])
    return slide


def build_s09_keyword_search(prs, notes):
    slide = base_slide(prs, "PHASE 3  ·  SEARCH FIRST", "Finding the pieces — no AI needed", notes)

    q = card(slide, 0.9, 1.9, 7.6, 0.6, fill=PANEL, line_color=TEAL, radius=0.3)
    txt(slide, 1.2, 1.9, 7.0, 0.6, [("\u201cWhat are the library opening hours?\u201d", 13.5, WHITE, True, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    txt(slide, 0.9, 2.7, 5.0, 0.35, [("THE KB, AS ACTUALLY STORED — A TABLE OF CHUNKS", 10, MUTED, True, False, None)])

    rows = [
        ("A", "\u201cLibrary opening hours: 8am\u201310pm on term days.\u201d", "library · opening · hours", "3", TEAL, True),
        ("B", "\u201cEnrolment closes at the census date.\u201d", "—", "0", PANEL_LINE, False),
        ("C", "\u201cExam Period Appendix: library access 24/7 during exam period.\u201d", "library", "1", PANEL_LINE, False),
    ]
    for i, (cid, text, shared, score, line_color, winner) in enumerate(rows):
        y = 3.1 + i * 0.78
        card(slide, 0.9, y, 8.7, 0.66, fill=PANEL, line_color=line_color)
        badge = circle(slide, 1.32, y + 0.33, 0.4, fill=TEAL if winner else PANEL_LINE)
        centered_label(badge, cid, 12, BG if winner else MUTED)
        txt(slide, 1.65, y, 4.55, 0.66, [(text, 10.5, WHITE, False, False, 1.1)], anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, 6.3, y, 1.95, 0.66, [(shared, 9.5, MUTED, False, True, 1.1)], anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, 8.35, y, 1.1, 0.66, [(("score " + score + (" \u2713" if winner else "")), 11,
                                          TEAL if winner else MUTED, True, False, None)], anchor=MSO_ANCHOR.MIDDLE)

    side = card(slide, 9.85, 1.9, 2.85, 3.0, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 10.1, 1.9, 2.35, 3.0, [
        ("Count shared words.", 12, WHITE, True, False, 1.4),
        ("Rare words weigh extra.", 12, WHITE, False, False, 1.4),
        ("Highest score wins.", 12, WHITE, False, False, 1.55),
        ("How search engines worked for decades.", 11, MUTED, False, True, 1.35),
        ("No AI anywhere.", 12.5, TEAL, True, False, None),
    ], anchor=MSO_ANCHOR.MIDDLE)

    term_chip(slide, 0.9, 5.55, 8.7, 0.75,
              "Keyword search (BM25)", "score each chunk by word overlap with the question, rare words counting extra.")

    txt(slide, 0.9, 6.5, 11.6, 0.4,
        [("Wire this in — and the bot suddenly works surprisingly well. Fast, cheap… and usually right.", 13, TEAL, False, True, None)])
    return slide


def build_s10_bug1(prs, notes):
    slide = base_slide(prs, "PHASE 3  ·  IT BREAKS", "Right words, wrong rule", notes)
    failure_chip(slide, "\u2717 confidently wrong", w=2.9)

    q = card(slide, 4.1, 1.85, 8.35, 0.52, fill=PANEL, line_color=TEAL, radius=0.3)
    txt(slide, 4.4, 1.85, 7.85, 0.52,
        [("Week 3 of term. Student asks: \u201cCan I access the library 24/7?\u201d", 12.5, WHITE, True, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    card(slide, 0.9, 2.7, 5.65, 2.0, fill=PANEL, line_color=WRONG)
    txt(slide, 1.15, 2.85, 5.15, 0.4, [("Exam Period Appendix", 13, WRONG, True, False, None)])
    txt(slide, 1.15, 3.3, 5.15, 0.6, [("\u201c…library access 24/7 during the exam period…\u201d", 12, WHITE, False, True, 1.2)])
    txt(slide, 1.15, 4.1, 5.15, 0.45, [("shares: library · access · 24/7  →  score 3  →  wins \u2717", 11, WRONG, True, False, None)])

    card(slide, 6.8, 2.7, 5.65, 2.0, fill=PANEL, line_color=TEAL)
    txt(slide, 7.05, 2.85, 5.15, 0.4, [("General rule — governs this week", 13, TEAL, True, False, None)])
    txt(slide, 7.05, 3.3, 5.15, 0.6, [("\u201cLibrary opening hours: 8am\u201310pm on term days.\u201d", 12, WHITE, False, True, 1.2)])
    txt(slide, 7.05, 4.1, 5.15, 0.45, [("shares: library  →  score 1  →  left behind", 11, MUTED, True, False, None)])

    bot = card(slide, 0.9, 5.0, 5.65, 0.85, fill=PANEL, line_color=WRONG)
    txt(slide, 1.15, 5.0, 5.15, 0.85,
        [("MetMate: \u201cYes — the library is open 24/7.\u201d", 12.5, WHITE, True, False, 1.25),
         ("Confident. Wrong. Every word from a real document.", 10.5, WRONG, False, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    quiet = card(slide, 6.8, 5.0, 5.65, 0.85, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 7.05, 5.0, 5.15, 0.85,
        [("The quieter failure: \u201cenrolment\u201d scores ZERO against a chunk", 11, WHITE, False, False, 1.25),
         ("that says \u201cregistration\u201d. Synonyms are invisible.", 11, WHITE, False, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    txt(slide, 0.9, 6.2, 11.6, 0.6,
        [("Matching exact words has a name: LEXICAL matching. What we actually need: match the MEANING —", 13, AMBER, False, True, 1.3),
         ("SEMANTIC matching. That is Phase 4.", 13, AMBER, False, True, None)])
    return slide


def build_s11_math_with_words(prs, notes):
    slide = base_slide(prs, "PHASE 4  ·  SEARCH BY MEANING", "Math with words", notes)

    txt(slide, 0.9, 1.85, 11.6, 0.4,
        [("Give every word a number — then \u201chow similar?\u201d becomes \u201chow close?\u201d, plain subtraction.",
          13.5, MUTED, False, True, None)])

    panel = card(slide, 0.9, 2.45, 8.0, 2.9, fill=PANEL, line_color=PANEL_LINE)
    axis_y = 4.15
    def line_x(value):
        return 4.9 + value * 0.5
    straight_line(slide, 1.4, axis_y, 8.5, axis_y, MUTED, 1.5)
    for value in (-5, 0, 5, 6):
        x = line_x(value)
        straight_line(slide, x, axis_y - 0.07, x, axis_y + 0.07, MUTED, 1.5)
        txt(slide, x - 0.4, axis_y + 0.14, 0.8, 0.3, [(str(value), 10, MUTED, False, False, None)],
            align=PP_ALIGN.CENTER)

    circle(slide, line_x(-5), axis_y, 0.18, fill=WRONG)
    txt(slide, line_x(-5) - 1.0, axis_y - 0.55, 2.0, 0.35, [("withdrawal", 12, WRONG, True, False, None)],
        align=PP_ALIGN.CENTER)
    circle(slide, line_x(5), axis_y, 0.18, fill=TEAL)
    txt(slide, line_x(5) - 1.0, axis_y - 0.55, 2.0, 0.35, [("enrolment", 12, TEAL, True, False, None)],
        align=PP_ALIGN.CENTER)
    circle(slide, line_x(6), axis_y, 0.18, fill=TEAL)
    txt(slide, line_x(6) - 1.0, axis_y + 0.5, 2.0, 0.35, [("registration", 12, TEAL, True, False, None)],
        align=PP_ALIGN.CENTER)

    txt(slide, line_x(5) - 1.1, 2.75, 2.9, 0.35, [("1 apart — similar \u2713", 11, TEAL, False, True, None)])
    txt(slide, line_x(-2.5) - 1.5, 4.85, 3.0, 0.35, [("10 apart — opposite", 11, MUTED, False, True, None)],
        align=PP_ALIGN.CENTER)

    term_chip(slide, 9.2, 2.45, 3.5, 1.3, "Embedding",
              "writing words as numbers so that close numbers = similar meaning.", def_size=10.5)

    take = card(slide, 9.2, 3.95, 3.5, 1.4, fill=PANEL, line_color=AMBER)
    txt(slide, 9.45, 3.95, 3.0, 1.4,
        [("Take ONE thing home today: this concept.", 12, AMBER, True, False, 1.3),
         ("It unlocks half of the LLM world.", 11, WHITE, False, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    hook = card(slide, 0.9, 5.65, 11.8, 0.95, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 5.65, 11.2, 0.95,
        [("But where does \u201cpass\u201d go on this line? Not similar to enrolment, not opposite either…", 13, WHITE, False, True, 1.25),
         ("one number per word is not enough.", 13, WHITE, False, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)
    return slide


def build_s12_vectors(prs, notes):
    slide = base_slide(prs, "PHASE 4  ·  SEARCH BY MEANING", "One number isn't enough — vectors", notes)

    panel = card(slide, 0.9, 1.95, 6.3, 4.55, fill=PANEL, line_color=PANEL_LINE)
    cx, cy = 3.75, 4.25
    sx, sy = 0.42, 0.38

    def px(value):
        return cx + value * sx

    def py(value):
        return cy - value * sy

    straight_line(slide, px(-5.6), cy, px(6.6), cy, MUTED, 1.25)
    straight_line(slide, cx, py(5.6), cx, py(-5.6), MUTED, 1.25)
    txt(slide, 1.15, 2.15, 2.45, 0.75, [
        ("x: withdrawal \u2194 enrolment", 10, MUTED, False, True, 1.4),
        ("y: fail \u2194 pass", 10, MUTED, False, True, None),
    ])

    points = [
        ("enrolment [5, 0]", 5, 0, TEAL, -1.55, -0.42),
        ("registration [6, 0]", 6, 0, TEAL, -0.6, 0.26),
        ("withdrawal [\u22125, 0]", -5, 0, TEAL, -0.85, -0.42),
        ("pass [0, 5]", 0, 5, AMBER, 0.14, -0.1),
        ("fail [0, \u22125]", 0, -5, AMBER, 0.14, -0.12),
    ]
    for label, vx, vy, color, dx, dy in points:
        circle(slide, px(vx), py(vy), 0.16, fill=color)
        txt(slide, px(vx) + dx, py(vy) + dy, 2.1, 0.3, [(label, 10.5, color, True, False, None)])

    course_ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px(2.5) - 0.17), Inches(py(2.5) - 0.17),
                                         Inches(0.34), Inches(0.34))
    course_ring.fill.background()
    outline(course_ring, AMBER, 2.0)
    circle(slide, px(2.5), py(2.5), 0.16, fill=WHITE)
    txt(slide, px(2.5) + 0.2, py(2.5) - 0.35, 2.2, 0.3, [("course [2.5, 2.5]", 11.5, WHITE, True, False, None)])

    term_chip(slide, 7.5, 1.95, 4.95, 0.9, "Vector",
              "a list of numbers. Here: [enrolment-ness, pass-ness].")
    term_chip(slide, 7.5, 3.0, 4.95, 0.9, "Dimension",
              "one slot in that list. A 3-number vector = 3-dimensional.")

    reuse = card(slide, 7.5, 4.05, 4.95, 1.3, fill=PANEL, line_color=AMBER)
    txt(slide, 7.75, 4.05, 4.45, 1.3,
        [("\u201ccourse\u201d needs NO new axis: [2.5, 2.5] — a bit of both.", 12, WHITE, True, False, 1.3),
         ("A few hundred dimensions can place ALL the words.", 11.5, MUTED, False, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    real = card(slide, 7.5, 5.5, 4.95, 1.0, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 7.75, 5.5, 4.45, 1.0,
        [("Production embeddings: ~300–3,000 dimensions.", 11.5, WHITE, False, False, 1.3),
         ("Same idea — just more room.", 11, MUTED, False, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    txt(slide, 0.9, 6.6, 6.3, 0.35,
        [("similar meaning = nearby points — distance IS similarity", 11.5, TEAL, False, True, None)],
        align=PP_ALIGN.CENTER)
    return slide


def build_s13_learned_numbers(prs, notes):
    slide = base_slide(prs, "PHASE 4  ·  SEARCH BY MEANING", "Who assigns the numbers?", notes)

    txt(slide, 0.9, 1.85, 11.6, 0.4,
        [("Typing numbers for every word by hand would take forever. They are LEARNED.", 13.5, MUTED, False, True, None)])

    cloze = card(slide, 0.9, 2.5, 6.6, 1.15, fill=PANEL, line_color=TEAL)
    txt(slide, 1.2, 2.65, 6.0, 0.5, [("\u201cStudents must  ______  before the census date.\u201d", 15, WHITE, True, False, None)])
    txt(slide, 1.2, 3.2, 6.0, 0.35, [("hide a word from a real sentence — make the computer guess it", 10.5, MUTED, False, True, None)])

    guesses = card(slide, 0.9, 3.95, 6.6, 2.0, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 4.1, 6.0, 0.35, [("the computer's guesses:", 11, MUTED, True, False, None)])
    bars = [
        ("enrol", 2.9, "62%", TEAL),
        ("register", 1.45, "31%", TEAL),
        ("party", 0.08, "0.1%", PANEL_LINE),
    ]
    for i, (word, bar_w, pct, color) in enumerate(bars):
        y = 4.55 + i * 0.45
        txt(slide, 1.2, y, 1.3, 0.35, [(word, 12, WHITE, True, False, None)])
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.6), Inches(y + 0.06), Inches(bar_w), Inches(0.22))
        solid(bar, color)
        no_line(bar)
        bar.adjustments[0] = 0.5
        txt(slide, 2.7 + bar_w, y, 0.9, 0.35, [(pct, 10.5, MUTED, False, False, None)])

    side = card(slide, 7.9, 2.5, 4.55, 3.45, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 8.2, 2.72, 4.0, 3.1, [
        ("No answer sheet needed —", 12.5, TEAL, True, False, 1.3),
        ("the hidden word IS the answer.", 12.5, TEAL, True, False, 1.5),
        ("Words that fit the same gaps (enrol, register) get pushed to nearby numbers.", 12, WHITE, False, False, 1.3),
        ("Words that never fit drift away.", 12, WHITE, False, False, 1.5),
        ("Millions of sentences later, the numbers have assigned themselves.", 12, MUTED, False, False, 1.3),
    ])

    term_chip(slide, 0.9, 6.15, 11.55, 0.75,
              "Representation learning",
              "the numbers are learned from data, not designed by hand. Look up later: king \u2212 man + woman \u2248 queen.")
    return slide


def build_s14_meaning_search(prs, notes):
    slide = base_slide(prs, "PHASE 4  ·  SEARCH BY MEANING", "MetMate searches by meaning", notes)

    table = card(slide, 0.9, 1.9, 6.45, 2.6, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.15, 2.05, 6.0, 0.35, [("THE STORED KB — every chunk now carries its vector", 10, MUTED, True, False, None)])
    kb_rows = [
        ("\u201cEnrolment corrections must be lodged before census.\u201d", "[ 5.2,  0.1, …]", TEAL),
        ("\u201cLibrary opening hours: 8am\u201310pm on term days.\u201d", "[\u22120.2,  4.7, …]", MUTED),
        ("\u201cWithdrawal without penalty: apply before week 4.\u201d", "[\u22124.8,  0.3, …]", MUTED),
    ]
    for i, (chunk_text, vector_text, color) in enumerate(kb_rows):
        y = 2.5 + i * 0.62
        txt(slide, 1.15, y, 4.1, 0.6, [(chunk_text, 10, WHITE, False, False, 1.1)], anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, 5.35, y, 1.85, 0.6, [(vector_text, 10.5, color, True, False, None)],
            anchor=MSO_ANCHOR.MIDDLE, font=MONO)

    term_chip(slide, 7.6, 1.9, 4.85, 0.95, "Similarity",
              "the closeness of two vectors. Official name: cosine similarity.")

    trap = card(slide, 7.6, 3.0, 4.85, 1.5, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 7.85, 3.0, 4.35, 1.5,
        [("And the library trap? A \u201c24/7\u201d question now retrieves BOTH library chunks — rule and appendix "
          "are both ABOUT library access. The model sees the full picture and answers with the condition. \u2713",
          10.5, WHITE, False, False, 1.25)], anchor=MSO_ANCHOR.MIDDLE)

    qb = card(slide, 0.9, 4.85, 4.5, 0.8, fill=PANEL, line_color=TEAL, radius=0.3)
    txt(slide, 1.15, 4.85, 4.0, 0.8, [("\u201cHow do I fix a mistake in my registration?\u201d", 11.5, WHITE, True, True, 1.2)],
        anchor=MSO_ANCHOR.MIDDLE)
    arrow_right(slide, 5.5, 5.25, 0.45, fill=TEAL)
    vec = card(slide, 6.05, 4.85, 1.85, 0.8, fill=PANEL, line_color=AMBER)
    txt(slide, 6.2, 4.85, 1.55, 0.8, [("[ 5.6,  0.2, …]", 11, AMBER, True, False, None)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    arrow_right(slide, 8.0, 5.25, 0.45, fill=TEAL)
    win = card(slide, 8.55, 4.75, 3.9, 1.0, fill=PANEL, line_color=TEAL)
    txt(slide, 8.8, 4.75, 3.4, 1.0,
        [("nearest chunk: \u201cEnrolment corrections…\u201d", 11, WHITE, True, False, 1.2),
         ("ZERO shared words — found by meaning \u2713", 10.5, TEAL, True, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    txt(slide, 0.9, 6.15, 11.6, 0.45,
        [("Same skeleton — one box upgraded. The bot just learned MEANING.", 14, AMBER, True, False, None)])
    txt(slide, 0.9, 6.6, 11.6, 0.35,
        [("Search with \u201cregistration\u201d — texts about \u201cenrolment\u201d now come to you.", 11.5, MUTED, False, True, None)])
    return slide


def build_s15_rag(prs, notes):
    slide = base_slide(prs, "ALL TOGETHER", "The whole picture — you built RAG", notes)

    left = 0.9
    right = 12.43

    txt(slide, left, 1.82, 6.0, 0.32, [("BUILT ONCE — BEFORE ANY QUESTION", 10, MUTED, True, False, None)])
    legend_items = [("phase 2", PANEL_LINE), ("phase 3", TEAL), ("phase 4", AMBER)]
    lx = 8.6
    for label, color in legend_items:
        swatch = card(slide, lx, 1.84, 0.34, 0.2, fill=PANEL, line_color=color, radius=0.25)
        txt(slide, lx + 0.42, 1.78, 0.95, 0.3, [(label, 10, MUTED, False, False, None)])
        lx += 1.45

    top_positions = flow_row(slide, left, right, 2.2, 0.8, [
        ("Documents (KB)", PANEL_LINE),
        ("Split into chunks", TEAL),
        ("Chunk \u2192 vector", AMBER),
        ("Store: chunk + vector", AMBER),
    ], size=11, arrow_color=PANEL_LINE)

    txt(slide, left, 3.28, 6.0, 0.32, [("EVERY QUESTION, LIVE", 10, MUTED, True, False, None)])
    bottom_positions = flow_row(slide, left, right, 3.62, 0.95, [
        ("Question", PANEL_LINE),
        ("Question \u2192 vector", AMBER),
        ("Find nearest chunks", TEAL),
        ("Build prompt:\nchunks + question", PANEL_LINE),
        ("Model writes", PANEL_LINE),
        ("Answer", PANEL_LINE),
    ], size=10, arrow_color=PANEL_LINE)

    store_x, store_w = top_positions[3]
    near_x, near_w = bottom_positions[2]
    straight_line(slide, store_x + store_w / 2, 3.0, near_x + near_w / 2, 3.62, MUTED, 1.25)

    model_x, model_w = bottom_positions[4]
    ai_badge = card(slide, model_x + model_w / 2 - 0.75, 4.62, 1.5, 0.32, fill=AMBER, line_color=AMBER, radius=0.5)
    centered_label(ai_badge, "the only AI box", 8.5, BG)

    def bracket(x1, x2, label):
        y = 5.15
        straight_line(slide, x1, y, x2, y, AMBER, 1.75)
        straight_line(slide, x1, y - 0.07, x1, y + 0.07, AMBER, 1.75)
        straight_line(slide, x2, y - 0.07, x2, y + 0.07, AMBER, 1.75)
        txt(slide, x1, y + 0.05, x2 - x1, 0.32, [(label, 10.5, AMBER, True, False, None)], align=PP_ALIGN.CENTER)

    qv_x, qv_w = bottom_positions[1]
    prompt_x, prompt_w = bottom_positions[3]
    answer_x, answer_w = bottom_positions[5]
    bracket(qv_x, near_x + near_w, "RETRIEVAL — find it")
    bracket(prompt_x, prompt_x + prompt_w, "AUGMENTED — add it")
    bracket(model_x, model_x + model_w, "GENERATION — write it")

    reveal = card(slide, 0.9, 5.75, 11.55, 0.9, fill=PANEL, line_color=TEAL)
    txt(slide, 1.2, 5.75, 11.0, 0.9,
        [("This recipe has a name: RAG — Retrieval-Augmented Generation.", 15, WHITE, True, False, 1.3),
         ("You didn't memorise it. You built it, failure by failure.", 12, TEAL, False, True, None)],
        anchor=MSO_ANCHOR.MIDDLE)
    return slide


def build_s16_real_world(prs, notes):
    slide = base_slide(prs, "PHASE 5  ·  TOWARD THE REAL WORLD", "What real chatbots add", notes)

    txt(slide, 0.9, 1.85, 11.6, 0.45,
        [("Between today's bot and a production bot: refinements, not new ideas. One line each — look them up when you need them.",
          12.5, MUTED, False, True, None)])

    concepts = [
        ("Hierarchical chunking", "store small chunks AND whole sections — detail questions and overview questions both work"),
        ("Context-aware chunking", "cut at natural boundaries — headings, clauses — not blindly every N words"),
        ("Hybrid search", "keyword + meaning search together; each catches what the other misses"),
        ("Reranking", "a second, more careful pass re-orders the top chunks before the prompt"),
        ("Conversation memory", "\u201cwhat about postgrads?\u201d only works with the chat so far — rewrite it into a full question"),
        ("Verification", "check every claim against the sources before showing it — if support is missing, say \u201cI don't know\u201d"),
    ]
    left = 0.9
    cols = 3
    gap = 0.25
    w = (SLIDE_W - 0.9 - left - gap * (cols - 1)) / cols
    h = 1.5
    for i, (term, gloss) in enumerate(concepts):
        row = i // cols
        col = i % cols
        x = left + col * (w + gap)
        y = 2.55 + row * (h + 0.3)
        card(slide, x, y, w, h, fill=PANEL, line_color=PANEL_LINE)
        txt(slide, x + 0.22, y + 0.15, w - 0.44, 0.4, [(term, 13.5, TEAL, True, False, None)])
        txt(slide, x + 0.22, y + 0.6, w - 0.44, 0.85, [(gloss, 10.5, MUTED, False, False, 1.2)])

    txt(slide, 0.9, 6.35, 11.6, 0.45,
        [("Every one of these earned its place the same way everything today did: a real failure demanded it.",
          12.5, AMBER, False, True, None)], align=PP_ALIGN.CENTER)
    return slide


def build_s17_takeaways(prs, notes):
    slide = base_slide(prs, "TAKE-AWAYS", "Three things to take home", notes)

    takeaways = [
        ("An AI product is mostly ordinary software", "one AI box — and YOUR code decides what it reads."),
        ("Embeddings", "words as numbers, meaning as distance. The one concept to keep."),
        ("RAG — Retrieve, Augment, Generate", "you built it today, phase by phase, failure by failure."),
    ]
    for i, (head, body) in enumerate(takeaways):
        y = 2.15 + i * 1.15
        card(slide, 0.9, y, 11.55, 0.95, fill=PANEL, line_color=TEAL if i == 1 else PANEL_LINE)
        num = circle(slide, 1.5, y + 0.475, 0.5, fill=AMBER if i == 1 else TEAL)
        centered_label(num, str(i + 1), 15, BG)
        box = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(10.2), Inches(0.95))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = head
        r1.font.name = FONT
        r1.font.size = Pt(15)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = "  —  " + body
        r2.font.name = FONT
        r2.font.size = Pt(13)
        r2.font.color.rgb = MUTED

    txt(slide, 0.9, 5.75, 11.55, 0.45,
        [("The recipe is public. Fitting it to YOUR domain is the part that's yours.", 13.5, AMBER, False, True, None)],
        align=PP_ALIGN.CENTER)
    txt(slide, 0.9, 6.25, 11.55, 0.6, [("Thank you — questions?", 24, WHITE, True, False, None)],
        align=PP_ALIGN.CENTER)
    return slide


NOTES = {
    "s01": (
        "[~1 min] Good afternoon everyone, thank you for having me. I'm Arif — Analyst Programmer at "
        "Wimmera CMA in regional Victoria. My day job is building software and mapping systems, and over "
        "the last year I have built several AI tools for my organisation. One of them is a chatbot that "
        "answers questions from official documents, and it is in real daily use. Today I will show you how "
        "such a chatbot really works — by building one from zero, in five phases. Each phase will hit a "
        "real problem, and the problem will force the next idea. I can't show my workplace's internal "
        "system, so we will use a stand-in: MetMate, an invented chatbot for Sydney Met itself, answering "
        "student questions about enrolment, the library, and exams. Same ideas, public-friendly example. "
        "One promise before we start: concepts over terminology. Every technical term will arrive only "
        "when the story needs it, and I will define it in plain English when it does. Nothing to memorise "
        "up front — and everything can be looked up later."
    ),
    "s02": (
        "[~2 min] Phase one. The simplest possible way to talk to your documents — and it needs no code at "
        "all. Open ChatGPT, attach the student handbook PDF, type your question. Done. What happens under "
        "the hood? For text, you can imagine it very simply: the full text of the document and your "
        "question are glued together into ONE long piece of text, and the model reads all of it and writes "
        "an answer. That's the picture — one big string of text. Keep it in your head, because the whole "
        "lecture is about improving this one picture. Two small notes. First, today we care about text "
        "only — what happens with images or complex files is a story for another day. Second, and "
        "honestly: this works. If you are one person with one document, this is genuinely the right "
        "solution, and my honest advice is: just do this. The rest of this lecture exists because a "
        "university is not one person."
    ),
    "s03": (
        "[~2.5 min] So why can't Sydney Met just tell 30,000 students 'go upload the handbook to "
        "ChatGPT'? Three problems. One — ChatGPT is somebody else's product. We can't put ChatGPT itself "
        "on the university website; to offer that experience ourselves we would have to build our own "
        "ChatGPT-like product, which is a huge job. Two — the alternative is that every student needs "
        "their own paid subscription. Thousands of seats, just to ask about the handbook. Three — and "
        "this is the quiet killer — a real answer often needs several documents at once: the handbook, "
        "the library rules, the exam procedures. And these documents evolve; they change through the "
        "year. Expecting every student to re-upload the right, current pile of PDFs into every private "
        "chat, every time — that will simply never happen. Now our first technical term of the day. That "
        "pile of official documents the bot must answer from has a name: the Knowledge Base — the KB. "
        "And notice the teal box on the slide — every time you see a box like that today, it means a new "
        "term just entered our story, defined in plain English, ready to be looked up later. So here is "
        "what we actually want: ONE bot, OUR OWN, sitting on the university website, always reading the "
        "CURRENT documents. That means we build. Phase two."
    ),
    "s04": (
        "[~2.5 min] Phase two: our own bot, the smallest one possible. Three pieces. A web page with a "
        "chat box — ordinary web development, many of you can build this already. A server behind it — "
        "ordinary backend code. And the model? We don't have our own model — but the maker of ChatGPT "
        "sells access to the machine behind it. That is the OpenAI API. A quick word on names, because "
        "it confuses everyone: ChatGPT is the product, the website you chat with. OpenAI is the company, "
        "and the API is how our program talks to their model directly. API — a doorway for programs: our "
        "server sends text over the internet, the model's reply comes back as text. No chat window "
        "involved. I should also say: you don't have to use a paid API at all. You can run a free model "
        "on your own machine — a popular tool for that is called Ollama — and everything in this lecture "
        "would work the same way. For a concrete picture, for the rest of the talk, let's say we use the "
        "OpenAI API. And what does our server actually send? Exactly the phase-one trick, just through "
        "pipes: one big string — every document in the KB, plus the student's question."
    ),
    "s05": (
        "[~2.5 min] Let's make that string concrete, because this exact shape carries the rest of the "
        "lecture. You can imagine the model receives a text like this: the word 'Context', then the "
        "documents pasted in; the word 'Question', then the student's question; then one instruction — "
        "'answer using only the context above'. Two terms arrive here. This whole message, everything we "
        "send, is called the PROMPT. And the reference material we paste in for the model to answer from "
        "is called the CONTEXT. Now the important realisation: this is ALL the model sees. It has no "
        "memory of Sydney Met, no hidden database, no idea our university exists. If the right fact is "
        "not somewhere in that prompt, the model can only guess. So the whole craft of a document chatbot "
        "comes down to one question: WHAT goes into this text? Hold that — it is the rest of the lecture."
    ),
    "s06": (
        "[~3 min] Before we go further, pause and look at what we have — because there is an important "
        "observation here. Any 'AI system' you meet is not fully AI. The web page — ordinary code. The "
        "server — ordinary code. Gluing strings together — ordinary code. Exactly ONE box in this "
        "picture is AI: the model. There is a clean way to tell the two worlds apart — not by "
        "definition, but by behaviour. Ordinary code is DETERMINISTIC: same input, same output, every "
        "single time. You can test it, trust it, debug it — the skills you already have from software "
        "engineering apply directly. The model is NON-DETERMINISTIC: ask the exact same question twice "
        "and the wording of the answer may differ. So when you hear 'AI system', picture this slide: one "
        "non-deterministic box, surrounded by ordinary, deterministic software. And here is the part I "
        "want you to remember as future developers: most of the work — and most of the value you add — "
        "lives in the deterministic part. That is where we spend the rest of today."
    ),
    "s07": (
        "[~2.5 min] Now phase two breaks. A university KB is not one handbook — it's hundreds of PDFs. "
        "Say 300 PDFs, roughly two million words. And our design sends ALL of it, for EVERY question — "
        "even for 'what are the library hours?'. I can tell you from experience what happens, because my "
        "first real bot was built exactly this way: it took about a MINUTE to answer. Nobody sits in a "
        "chat window waiting a minute. And think about what we are doing — almost everything we send is "
        "completely irrelevant to the question being asked. One honest disclaimer before we fix it: we "
        "will NOT open up how the model itself works today. That is a deep field of its own, and for "
        "this lecture the model stays a very capable magic box. Our job — the chatbot developer's job — "
        "is to feed that box well: fast, and accurate. So the direction is obvious: send LESS. But which "
        "part? That question is phase three."
    ),
    "s08": (
        "[~3 min] Here's the idea. For any single question, only a few paragraphs out of the whole KB "
        "actually matter. So — ahead of time — we cut every document into small pieces. New term: a "
        "CHUNK. A chunk is a small passage cut from a document, a few paragraphs, ideally about one "
        "topic. Then, when a question arrives, we pick the top five or ten chunks that look most "
        "relevant, and we build exactly the same prompt as before — same Context-Question shape — just "
        "with those few chunks as the context instead of everything. Two hundred pages become two pages, "
        "per question. And one framing I really want you to keep. How the model answers from the text we "
        "give it — that is not our department; we take it as given. But choosing WHAT it reads — that is "
        "precisely where you, the chatbot developer, show your skill. Answering like ChatGPT does is a "
        "giant AI problem; picking the right text is an engineering problem, and it is YOUR problem. "
        "Almost everything that made my real bot better over time happened in this choosing step."
    ),
    "s09": (
        "[~3.5 min] So how do we pick the right chunks? First question to ask: do we need AI for this? "
        "And the honest answer is NO — and I want you to see that concretely. Here is the KB as it is "
        "actually stored: a plain table of chunks. Nothing exotic. The question comes in: 'What are the "
        "library opening hours?' Now just count shared words. Chunk A — the library hours rule — shares "
        "three: library, opening, hours. Chunk B, about enrolment — zero. Chunk C, the exam appendix — "
        "one: library. Chunk A wins, and chunk A is genuinely the right passage. Add one refinement — "
        "rare words count extra, because 'library' says more than 'the' — and you have essentially "
        "rebuilt BM25, the scoring recipe search engines ran on for decades. This is how search worked "
        "before AI — and notice, it is still ordinary deterministic code: a word table and counting. "
        "Build this, wire it into our pipeline, and something surprising happens: the bot suddenly feels "
        "GOOD. Fast, cheap, and usually right. My real bot at this stage was already impressing people. "
        "Usually right. Remember that word — usually."
    ),
    "s10": (
        "[~3.5 min] Here is the failure that ends phase three — and it is a real one; my bot did exactly "
        "this, and this example is its public stand-in. Week three of term. A student asks: 'Can I "
        "access the library 24/7?' Count the words. The Exam Period Appendix contains 'library', "
        "'access', AND '24/7' — score three, top result. The general rule — the one that actually "
        "governs an ordinary week — says 'library opening hours 8am to 10pm on term days'. It shares "
        "one word. Score one. Left behind. So MetMate answers, fluently and confidently: 'Yes — the "
        "library is open 24/7.' Wrong. And notice the scary part: every single word of that answer came "
        "from a real official document. Word counting rewards the chunk that ECHOES the question — not "
        "the rule that APPLIES. There's a quieter failure hiding underneath too: a question about "
        "'enrolment' scores ZERO against a chunk that says 'registration'. Synonyms are invisible to "
        "word counting. Both failures point the same way. Matching exact words — the technical name is "
        "LEXICAL matching — is not enough. We need to match MEANING — SEMANTIC matching. How on earth "
        "does a computer compute with meaning? That is phase four, and it is the heart of this lecture."
    ),
    "s11": (
        "[~3 min] To match meaning, we need to be able to do MATH with words. Take three words from our "
        "university world: enrolment, registration, withdrawal. You and I know enrolment and "
        "registration are nearly the same thing, and withdrawal is roughly the opposite. The question "
        "is: how can a computer measure that? Here is the simple, almost silly idea: give every word a "
        "NUMBER. Enrolment: 5. Registration: 6. Withdrawal: minus 5. Now 'how similar?' becomes 'how "
        "close?' — plain subtraction. Five and six are one apart: close, similar. Five and minus five "
        "are ten apart: far, opposite. This idea — writing words as numbers so that close numbers mean "
        "similar meaning — is called an EMBEDDING. And I will say this directly: if you take only ONE "
        "thing home from this lecture, I would be very glad if it is this concept. Understand embeddings "
        "and half of the modern LLM world becomes accessible to you. But there's a problem with one "
        "number per word. Along comes the word 'pass'. Where does it go on this line? It's not similar "
        "to enrolment. It's not the opposite either. It's just… about something else. One number is not "
        "enough."
    ),
    "s12": (
        "[~4 min] The fix: give each word SEVERAL numbers instead of one — a list. New term: a VECTOR. "
        "In our context, a vector is just a list of numbers. Let the first number measure the "
        "enrolment-versus-withdrawal direction, and the second measure the pass-versus-fail direction. "
        "Now watch. Enrolment: [5, 0] — strongly about enrolling, nothing to do with passing or "
        "failing. Registration: [6, 0]. Withdrawal: [minus 5, 0]. Pass: [0, 5] — nothing to do with "
        "enrolment. And fail: [0, minus 5], a bonus word. Each slot in the list is called a DIMENSION — "
        "a 3-number vector is 3-dimensional, and so on. Now think about the word 'course'. It's related "
        "to enrolment — you enrol in a course. It's related to passing — you pass a course. But it is "
        "not quite either. Do we need a third slot for it? We could add one — or we can reuse what we "
        "already have: course = [2.5, 2.5]. A bit of both. And that is the deep trick of embeddings: we "
        "do NOT need a new dimension for every word. A few hundred well-chosen dimensions can place "
        "hundreds of thousands of words, each word a point, similar words nearby. Real production "
        "embeddings use roughly three hundred to three thousand dimensions — not two — but it is "
        "exactly this idea, just with more room. Meaning has become geometry: distance IS similarity."
    ),
    "s13": (
        "[~2.5 min] One question should be bothering you: who assigns all these numbers? If we typed "
        "them by hand it would literally take forever — and we'd argue for a week about 'course'. The "
        "answer: nobody assigns them. They are LEARNED. Here is one classic recipe. Take millions of "
        "real sentences. Hide a word: 'Students must ______ before the census date.' Make the computer "
        "guess the missing word. Notice the beautiful part — this needs no human answer sheet, because "
        "the hidden word IS the answer. Words that fit the same gaps — enrol, register — get pushed "
        "toward nearby numbers. Words that never fit that gap — 'party' — drift away. Repeat over "
        "millions of sentences, and the numbers assign themselves. The general name for this — learning "
        "the number-form of things from data instead of designing it by hand — is REPRESENTATION "
        "LEARNING. And there is a famous party trick you can look up later: with vectors learned this "
        "way, king minus man plus woman lands almost exactly on queen. The arithmetic works because "
        "meaning became geometry."
    ),
    "s14": (
        "[~3.5 min] Now plug embeddings into MetMate. In the build-once lane, when we cut the KB into "
        "chunks, we now also compute each chunk's vector and store both — you can see the stored KB is "
        "still just a table: chunk text, plus its vector. In the live lane, the question becomes a "
        "vector too, and 'search' now means: find the chunks whose vectors sit CLOSEST to the "
        "question's vector. The closeness score between two vectors has an official name — cosine "
        "similarity — one more look-up-able term; it is nothing more than 'how close are these two "
        "lists of numbers'. Watch it work. 'How do I fix a mistake in my registration?' The right "
        "chunk says 'enrolment corrections' — ZERO shared words with the question; keyword search "
        "scores it zero and never finds it. But its vector sits right next to the question's vector — "
        "found by meaning. Search with 'registration', and texts about 'enrolment' now come to you. "
        "And our library trap? Ask about 24/7 access and the bot now retrieves BOTH library chunks — "
        "the general rule and the appendix — because both are ABOUT library access times. The model "
        "sees the full picture and answers with the condition: 8 to 10 normally, 24/7 only during "
        "exams. Same skeleton, one box upgraded — the bot just learned meaning."
    ),
    "s15": (
        "[~3 min] Zoom out and look at the whole machine — colour-coded by the phase that forced each "
        "box into existence. Grey: the phase-two skeleton — page, server, prompt, model, answer. Teal: "
        "phase three — cut into chunks, pick the best ones. Amber: phase four — vectors on both lanes. "
        "Nothing in this picture was designed up front; every box exists because a failure demanded "
        "it. Now the payoff. This exact recipe — SEARCH a knowledge base, PASTE the winners into the "
        "prompt, let the model WRITE — has a name in the industry: RAG. Retrieval-Augmented "
        "Generation. Retrieval — find it. Augmented — add it to the prompt. Generation — the model "
        "writes. You did not memorise RAG today; you BUILT it, and now you own it. Two honest "
        "footnotes. This diagram is ONE way to build a document chatbot — the field is young, there is "
        "no standard blueprint; this is the shape my real bot uses. And notice again: exactly one box "
        "in this whole picture is generative AI. Even the embedding step, once trained, is "
        "deterministic — same text in, same numbers out, every time. Mostly ordinary software, "
        "remember, around one remarkable box."
    ),
    "s16": (
        "[~1.5 min] Phase five, and I'll keep it to one line per idea — these are the refinements that "
        "separate today's teaching bot from a production bot, and every one of them is look-up-able. "
        "Hierarchical chunking: store small chunks AND whole sections, so detail questions and "
        "overview questions both work. Context-aware chunking: cut at natural boundaries — headings, "
        "clauses — not blindly every five hundred words. Hybrid search: run keyword search and meaning "
        "search together; each catches what the other misses. Reranking: a second, more careful pass "
        "re-orders the top chunks before they enter the prompt. Conversation memory: 'what about "
        "postgrads?' only makes sense given the chat so far — rewrite it into a full standalone "
        "question first. And verification: before showing an answer, check every claim against the "
        "source documents — and if the support is not there, say honestly, 'I don't know'. Each of "
        "these earned its place in my real bot the same way everything today did: a real failure "
        "demanded it."
    ),
    "s17": (
        "[~1.5 min] Three things to take home. One: an AI product is mostly ordinary software. One AI "
        "box — and your code, your engineering, decides what it reads. The skills you are already "
        "building apply directly. Two: embeddings. Words as numbers, meaning as distance. That is the "
        "one concept I asked you to keep — it unlocks half of what you will read about LLMs from here. "
        "Three: RAG — retrieve, augment, generate. You built it today, phase by phase, failure by "
        "failure — so it is yours now. The recipe is public; anyone can follow it. Making it work for "
        "a real domain — your future workplace's documents, rules, quirks — that part is not in any "
        "tutorial. That part is yours. Thank you — I'm happy to take questions."
    ),
}


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_s01_title(prs, NOTES["s01"])
    build_s02_upload(prs, NOTES["s02"])
    build_s03_doesnt_scale(prs, NOTES["s03"])
    build_s04_own_bot(prs, NOTES["s04"])
    build_s05_prompt(prs, NOTES["s05"])
    build_s06_ai_vs_not(prs, NOTES["s06"])
    build_s07_too_slow(prs, NOTES["s07"])
    build_s08_send_less(prs, NOTES["s08"])
    build_s09_keyword_search(prs, NOTES["s09"])
    build_s10_bug1(prs, NOTES["s10"])
    build_s11_math_with_words(prs, NOTES["s11"])
    build_s12_vectors(prs, NOTES["s12"])
    build_s13_learned_numbers(prs, NOTES["s13"])
    build_s14_meaning_search(prs, NOTES["s14"])
    build_s15_rag(prs, NOTES["s15"])
    build_s16_real_world(prs, NOTES["s16"])
    build_s17_takeaways(prs, NOTES["s17"])

    prs.save(OUT_PATH)
    print("saved", OUT_PATH)


if __name__ == "__main__":
    main()
