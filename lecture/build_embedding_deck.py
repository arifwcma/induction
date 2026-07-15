import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from build_deck import (
    BG, PANEL, PANEL_LINE, TEAL, AMBER, WHITE, MUTED, WRONG, FOOTER_BG,
    FONT, MONO, SLIDE_W, SLIDE_H,
    solid, no_line, outline, txt, card, circle, arrow_right, straight_line,
    base_slide, term_chip,
)


HERE = os.path.dirname(os.path.abspath(__file__))
OUT_PATH = os.path.join(HERE, "embedding_deck.pptx")

EYEBROW = "PHASE 4  ·  SEARCH BY MEANING"
TITLE = "Who assigns the numbers?"

SPACE_X = 7.05
SPACE_Y = 2.0
SPACE_W = 5.4
SPACE_H = 4.55
SPACE_CENTRE_X = SPACE_X + SPACE_W / 2
SPACE_CENTRE_Y = SPACE_Y + SPACE_H / 2 + 0.15
SPACE_SCALE_X = 0.33
SPACE_SCALE_Y = 0.30


def space_x(value):
    return SPACE_CENTRE_X + value * SPACE_SCALE_X


def space_y(value):
    return SPACE_CENTRE_Y - value * SPACE_SCALE_Y


def draw_space_panel(slide, heading):
    card(slide, SPACE_X, SPACE_Y, SPACE_W, SPACE_H, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, SPACE_X + 0.25, SPACE_Y + 0.15, SPACE_W - 0.5, 0.35,
        [(heading, 10.5, MUTED, True, False, None)])
    straight_line(slide, space_x(-5.6), SPACE_CENTRE_Y, space_x(5.6), SPACE_CENTRE_Y, PANEL_LINE, 1.0)
    straight_line(slide, SPACE_CENTRE_X, space_y(5.0), SPACE_CENTRE_X, space_y(-5.0), PANEL_LINE, 1.0)


def plot_word(slide, value_x, value_y, label, color, label_dx, label_dy):
    circle(slide, space_x(value_x), space_y(value_y), 0.16, fill=color)
    txt(slide, space_x(value_x) + label_dx, space_y(value_y) + label_dy, 2.0, 0.3,
        [(label, 11, color, True, False, None)])


def cloze_card(slide):
    card(slide, 0.9, 2.55, 5.6, 1.15, fill=PANEL, line_color=TEAL)
    txt(slide, 1.2, 2.72, 5.0, 0.5,
        [("\u201cStudents must  ______  before the census date.\u201d", 15, WHITE, True, False, None)])
    txt(slide, 1.2, 3.28, 5.0, 0.35,
        [("hide one real word — let the computer guess it", 10.5, MUTED, False, True, None)])


def guess_bars(slide, top_word_confident):
    card(slide, 0.9, 3.95, 5.6, 2.0, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 4.1, 5.0, 0.35, [("the computer's guess:", 11, MUTED, True, False, None)])
    if top_word_confident:
        bars = [("enrol", 2.4, "62%", TEAL), ("register", 1.2, "31%", TEAL), ("party", 0.08, "0.1%", PANEL_LINE)]
    else:
        bars = [("enrol", 0.85, "34%", MUTED), ("register", 0.78, "31%", MUTED), ("party", 0.72, "29%", MUTED)]
    for i, (word, bar_w, pct, color) in enumerate(bars):
        y = 4.55 + i * 0.45
        txt(slide, 1.2, y, 1.3, 0.35, [(word, 12, WHITE, True, False, None)])
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.6), Inches(y + 0.06), Inches(bar_w), Inches(0.22))
        solid(bar, color)
        no_line(bar)
        bar.adjustments[0] = 0.5
        txt(slide, 2.7 + bar_w, y, 0.9, 0.35, [(pct, 10.5, MUTED, False, False, None)])


def caption(slide, text, color=MUTED):
    txt(slide, 0.9, 1.85, 11.6, 0.4, [(text, 13.5, color, False, True, None)])


def build_f1_random(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES["f1"])
    caption(slide, "Typing a number for every word by hand is impossible. So every word STARTS as random numbers.")

    card(slide, 0.9, 2.55, 5.6, 3.4, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 2.75, 5.0, 0.4, [("each word begins as noise:", 11, MUTED, True, False, None)])
    starts = [("enrol", "[ -1.2,  0.1 ]"), ("register", "[  1.3, -0.2 ]"), ("party", "[ -0.9,  0.4 ]")]
    for i, (word, vector_text) in enumerate(starts):
        y = 3.35 + i * 0.75
        txt(slide, 1.2, y, 1.8, 0.5, [(word, 14, WHITE, True, False, None)], anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, 3.1, y, 3.0, 0.5, [(vector_text, 13, MUTED, True, False, None)], anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    txt(slide, 1.2, 5.55, 5.0, 0.35, [("meaningless — assigned by a coin toss", 10.5, WRONG, False, True, None)])

    draw_space_panel(slide, "THE NUMBER SPACE — starts as pure noise")
    plot_word(slide, -3.2, 2.6, "enrol", MUTED, 0.18, -0.3)
    plot_word(slide, 2.8, -2.8, "register", MUTED, 0.18, -0.1)
    plot_word(slide, 0.6, 3.4, "party", MUTED, 0.18, -0.3)
    return slide


def build_f2_game(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES["f2"])
    caption(slide, "The training game: take a real sentence, hide one word.")

    cloze_card(slide)
    hint = card(slide, 0.9, 3.95, 5.6, 2.0, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 4.15, 5.0, 1.6, [
        ("No human answer sheet needed.", 13, TEAL, True, False, 1.35),
        ("The hidden word IS the answer.", 13, WHITE, False, False, 1.35),
        ("We have millions of real sentences — so millions of free practice questions.", 12, MUTED, False, False, 1.3),
    ])

    draw_space_panel(slide, "THE NUMBER SPACE — still noise")
    plot_word(slide, -3.2, 2.6, "enrol", MUTED, 0.18, -0.3)
    plot_word(slide, 2.8, -2.8, "register", MUTED, 0.18, -0.1)
    plot_word(slide, 0.6, 3.4, "party", MUTED, 0.18, -0.3)
    return slide


def build_f3_guess(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES["f3"])
    caption(slide, "The computer guesses the blank from the words' CURRENT numbers. Random numbers \u2192 a random guess.")

    cloze_card(slide)
    guess_bars(slide, top_word_confident=False)

    draw_space_panel(slide, "THE NUMBER SPACE — still noise")
    plot_word(slide, -3.2, 2.6, "enrol", MUTED, 0.18, -0.3)
    plot_word(slide, 2.8, -2.8, "register", MUTED, 0.18, -0.1)
    plot_word(slide, 0.6, 3.4, "party", MUTED, 0.18, -0.3)
    txt(slide, SPACE_X + 0.25, SPACE_Y + SPACE_H - 0.55, SPACE_W - 0.5, 0.4,
        [("the guess is bad — because the numbers are still random", 10.5, MUTED, False, True, None)],
        align=PP_ALIGN.CENTER)
    return slide


def build_f4_nudge(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES["f4"])
    caption(slide, "Reveal the true word. Nudge the numbers a little so next time the guess is better.")

    cloze_card(slide)
    guess_bars(slide, top_word_confident=False)

    reveal = card(slide, 0.9, 6.05, 5.6, 0.7, fill=PANEL, line_color=TEAL)
    txt(slide, 1.15, 6.05, 5.1, 0.7,
        [("true word: enrol   \u2192   nudge enrol's numbers toward this context \u2713", 11.5, TEAL, True, False, None)],
        anchor=MSO_ANCHOR.MIDDLE)

    draw_space_panel(slide, "THE NUMBER SPACE — one tiny nudge")
    circle(slide, space_x(-3.2), space_y(2.6), 0.16, fill=PANEL_LINE)
    straight_line(slide, space_x(-3.2), space_y(2.6), space_x(-2.3), space_y(1.9), AMBER, 1.75)
    plot_word(slide, -2.3, 1.9, "enrol", AMBER, 0.18, -0.3)
    plot_word(slide, 2.8, -2.8, "register", MUTED, 0.18, -0.1)
    plot_word(slide, 0.6, 3.4, "party", MUTED, 0.18, -0.3)
    txt(slide, SPACE_X + 0.25, SPACE_Y + SPACE_H - 0.55, SPACE_W - 0.5, 0.4,
        [("one sentence moves it a hair. The magic is in the REPETITION.", 10.5, MUTED, False, True, None)],
        align=PP_ALIGN.CENTER)
    return slide


def build_f5_company(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES["f5"])
    caption(slide, "The key: words that fill the SAME blanks get nudged the SAME way \u2014 so they drift together.")

    panel = card(slide, 0.9, 2.55, 5.6, 3.4, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 2.7, 5.0, 0.35, [("both words fit these gaps:", 11, MUTED, True, False, None)])
    fits = [
        "\u201cStudents must ___ before the census date.\u201d",
        "\u201cYou can ___ online through the portal.\u201d",
        "\u201cThe deadline to ___ is week two.\u201d",
    ]
    for i, sentence in enumerate(fits):
        y = 3.15 + i * 0.55
        txt(slide, 1.2, y, 5.1, 0.5, [(sentence, 11.5, WHITE, False, True, 1.15)], anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, 1.2, 4.9, 5.1, 0.5,
        [("enrol and register keep swapping into the same slots \u2192 pulled close.", 11.5, TEAL, True, False, 1.2)])
    txt(slide, 1.2, 5.5, 5.1, 0.4,
        [("party never fits \u2192 pushed away.", 11.5, WRONG, True, False, None)])

    draw_space_panel(slide, "THE NUMBER SPACE — order emerging")
    plot_word(slide, 3.5, 0.8, "enrol", TEAL, 0.2, -0.32)
    plot_word(slide, 4.6, -0.2, "register", TEAL, 0.2, 0.08)
    straight_line(slide, space_x(-3.2), space_y(2.6), space_x(3.5), space_y(0.8), AMBER, 1.0)
    straight_line(slide, space_x(2.8), space_y(-2.8), space_x(4.6), space_y(-0.2), AMBER, 1.0)
    plot_word(slide, -3.5, 4.2, "party", WRONG, 0.2, -0.32)
    straight_line(slide, space_x(0.6), space_y(3.4), space_x(-3.5), space_y(4.2), WRONG, 1.0)
    return slide


def build_f6_map(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES["f6"])
    caption(slide, "Millions of sentences later, the noise has organised itself \u2014 the map from a few slides ago.")

    panel = card(slide, 0.9, 2.55, 5.6, 3.4, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 2.75, 5.0, 1.9, [
        ("Nobody placed these points.", 13.5, TEAL, True, False, 1.4),
        ("Similar words sit close because they kept the same company.", 12.5, WHITE, False, False, 1.35),
        ("Opposite words sit far apart.", 12.5, WHITE, False, False, 1.35),
        ("The numbers assigned themselves.", 12.5, AMBER, True, False, 1.35),
    ])
    txt(slide, 1.2, 5.55, 5.1, 0.35,
        [("this is exactly the enrolment / registration / withdrawal geometry", 10.5, MUTED, False, True, None)])

    draw_space_panel(slide, "THE NUMBER SPACE — meaning, learned")
    plot_word(slide, 4.4, 0.8, "enrol", TEAL, 0.2, -0.32)
    plot_word(slide, 5.2, 0.1, "register", TEAL, 0.2, 0.08)
    plot_word(slide, -4.6, 0.4, "withdrawal", WRONG, -0.35, -0.32)
    plot_word(slide, -3.2, 4.4, "party", MUTED, 0.2, -0.32)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(space_x(4.8) - 0.85), Inches(space_y(0.4) - 0.55),
                                  Inches(1.7), Inches(1.1))
    ring.fill.background()
    outline(ring, TEAL, 1.75)
    txt(slide, space_x(4.8) - 0.85, space_y(0.4) + 0.55, 1.7, 0.3,
        [("similar \u2014 close", 9.5, TEAL, True, False, None)], align=PP_ALIGN.CENTER)
    return slide


def build_f7_name(prs):
    slide = base_slide(prs, EYEBROW, TITLE, NOTES["f7"])
    caption(slide, "The numbers were learned from data, not designed by hand.")

    term_chip(slide, 0.9, 2.55, 11.55, 1.0,
              "Representation learning",
              "learning the number-form of things from raw data, instead of hand-designing it.")

    idea = card(slide, 0.9, 3.9, 5.6, 2.6, fill=PANEL, line_color=PANEL_LINE)
    txt(slide, 1.2, 4.1, 5.0, 2.2, [
        ("\u201cKnow a word by the company it keeps.\u201d", 13.5, WHITE, True, True, 1.35),
        ("", 8, MUTED, False, False, 1.0),
        ("Fill the same blanks \u2192 similar numbers.", 12, MUTED, False, False, 1.3),
        ("Never share a blank \u2192 far apart.", 12, MUTED, False, False, 1.3),
        ("No human ever labelled a single meaning.", 12, TEAL, True, False, 1.3),
    ])

    trick = card(slide, 6.85, 3.9, 5.6, 2.6, fill=PANEL, line_color=AMBER)
    txt(slide, 7.15, 4.15, 5.0, 0.4, [("a famous look-it-up-later party trick:", 11, AMBER, True, False, None)])
    txt(slide, 7.15, 4.75, 5.0, 0.8,
        [("king \u2212 man + woman  \u2248  queen", 20, WHITE, True, False, None)], font=MONO)
    txt(slide, 7.15, 5.7, 5.0, 0.7,
        [("the arithmetic works because meaning became geometry \u2014 the whole point of embeddings.",
          11.5, MUTED, False, True, 1.25)])
    return slide


NOTES = {
    "f1": (
        "One question should be bothering you: who assigns all these numbers? Typing them by hand is "
        "impossible - and we'd argue for a week about a single word. So we don't. Every word simply STARTS "
        "with random numbers - a coin toss. Right now they mean nothing: enrol, register and party are just "
        "three random dots, scattered anywhere. Watch what happens to this cloud of noise."
    ),
    "f2": (
        "Here is the game that fixes it. Take a real sentence from our documents and hide one word: "
        "'Students must ___ before the census date.' The beautiful part - this needs no human answer sheet, "
        "because the hidden word IS the answer. And we have millions of real sentences, so millions of free "
        "practice questions, for nothing."
    ),
    "f3": (
        "The computer tries to guess the blank, using the CURRENT numbers of the surrounding words. At the "
        "start those numbers are random, so the guess is basically a shrug - enrol, register, party all "
        "roughly equal. That's fine. A wrong guess is exactly what teaches it."
    ),
    "f4": (
        "Now we reveal the true word: enrol. The computer was not confident enough, so it nudges the numbers "
        "a tiny bit - shifting enrol's numbers toward this kind of context, so next time it leans more "
        "toward enrol. One sentence moves things by a hair. The magic is that we do this millions of times."
    ),
    "f5": (
        "Here is the heart of the whole idea. The words enrol and register fill the SAME blanks, over and "
        "over, across the whole language - 'you can ___ online', 'the deadline to ___'. Every time, both get "
        "nudged in the same direction. So they drift toward each other and end up with similar numbers. "
        "'party' never fits these gaps, so it is pushed away. Nobody told the computer they were synonyms - "
        "the shared company did."
    ),
    "f6": (
        "Fast-forward over millions of sentences and the random cloud has quietly organised itself. Similar "
        "words sit close because they kept the same company; opposites sit far apart. This is exactly the "
        "enrolment-registration-withdrawal map we drew earlier - except nobody placed a single point. The "
        "numbers assigned themselves."
    ),
    "f7": (
        "This whole trick - learning the number-form of words from raw data instead of designing it by hand "
        "- is called representation learning. The old linguists' line captures it: know a word by the "
        "company it keeps. And a famous party trick you can look up later: with vectors learned this way, "
        "king minus man plus woman lands almost exactly on queen. The arithmetic works because meaning has "
        "become geometry - which is the entire reason embeddings are so powerful."
    ),
}


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_f1_random(prs)
    build_f2_game(prs)
    build_f3_guess(prs)
    build_f4_nudge(prs)
    build_f5_company(prs)
    build_f6_map(prs)
    build_f7_name(prs)

    prs.save(OUT_PATH)
    print("saved", OUT_PATH)


if __name__ == "__main__":
    main()
