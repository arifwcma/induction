from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


BG = RGBColor(0x0B, 0x1B, 0x33)
PANEL = RGBColor(0x14, 0x28, 0x47)
PANEL_LINE = RGBColor(0x25, 0x3B, 0x5E)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
AMBER = RGBColor(0xF5, 0xB3, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9F, 0xB3, 0xC8)
FOOTER_BG = RGBColor(0x0F, 0x22, 0x40)

FONT = "Segoe UI"


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, size, color, bold, italic, spacing) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = FONT
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return box


def set_letterspacing(run, points):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(points * 100)))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
solid(bg, BG)
no_line(bg)

# faint accent bar top-left
strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), prs.slide_height)
solid(strip, TEAL)
no_line(strip)

# eyebrow
eb = txt(slide, 0.65, 0.42, 11.0, 0.4,
         [("GUEST LECTURE   ·   A REAL-WORLD AI PROJECT", 12.5, TEAL, True, False, None)])
set_letterspacing(eb.text_frame.paragraphs[0].runs[0], 2.2)

# title
txt(slide, 0.62, 0.78, 12.1, 1.0,
    [("Talk to Your Documents", 42, WHITE, True, False, 1.0)])

# accent underline
ul = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(1.72), Inches(1.7), Inches(0.07))
solid(ul, AMBER)
no_line(ul)

# subtitle
txt(slide, 0.66, 1.9, 12.0, 0.7,
    [("Building an AI assistant that answers from your PDFs — the journey from a naive bot to one you can actually trust.",
      17, MUTED, False, False, 1.1)])

# ---- 5 step cards ----
steps = [
    ("1", "Upload a PDF", "The easy way: drop it in ChatGPT.",
     "But — needs a paid seat, and won't scale to 1000s of docs."),
    ("2", "Teach it meaning", "Turn words into numbers (embeddings).",
     "queen − female + male ≈ king"),
    ("3", "Find the right passage", "Search by meaning, not keywords (RAG).",
     "But — it grabbed the wrong rule and answered confidently."),
    ("4", "Know what it knows", "Track every topic it actually covers.",
     "But — it refused a question it clearly had the answer to."),
    ("5", "Never make it up", "Check the answer against the source.",
     "The goal: confident only when it's genuinely right."),
]

left = 0.65
right = 13.333 - 0.65
usable = right - left
n = len(steps)
gap = 0.34
card_w = (usable - gap * (n - 1)) / n
card_y = 2.72
card_h = 2.72

for i, (num, title, desc, hint) in enumerate(steps):
    cx = left + i * (card_w + gap)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(cx), Inches(card_y), Inches(card_w), Inches(card_h))
    solid(card, PANEL)
    line(card, PANEL_LINE, 1.0)
    card.adjustments[0] = 0.06

    # number badge
    dia = 0.62
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(cx + 0.22), Inches(card_y + 0.24), Inches(dia), Inches(dia))
    solid(badge, TEAL)
    no_line(badge)
    btf = badge.text_frame
    btf.margin_left = 0; btf.margin_right = 0; btf.margin_top = 0; btf.margin_bottom = 0
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = num
    br.font.name = FONT; br.font.size = Pt(20); br.font.bold = True
    br.font.color.rgb = BG

    # title
    txt(slide, cx + 0.22, card_y + 1.02, card_w - 0.44, 0.7,
        [(title, 15.5, WHITE, True, False, 1.0)])
    # desc
    txt(slide, cx + 0.22, card_y + 1.62, card_w - 0.44, 0.6,
        [(desc, 10.5, MUTED, False, False, 1.05)])
    # hint (amber)
    txt(slide, cx + 0.22, card_y + 2.18, card_w - 0.44, 0.5,
        [(hint, 9.5, AMBER, False, True, 1.03)])

    # chevron between cards
    if i < n - 1:
        chev = txt(slide, cx + card_w - 0.02, card_y + 0.9, gap + 0.04, 0.5,
                   [("›", 26, TEAL, True, False, None)], align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE)

# ---- concept chips ----
chips = ["Embeddings", "Retrieval (RAG)", "Hallucination", "Reliability"]
chip_y = 5.72
chip_h = 0.52
chip_w = 2.7
chip_gap = 0.34
total = chip_w * len(chips) + chip_gap * (len(chips) - 1)
start = (13.333 - total) / 2
label = txt(slide, 0.65, chip_y - 0.02, 2.6, chip_h,
            [("YOU'LL MEET:", 11, MUTED, True, False, None)], anchor=MSO_ANCHOR.MIDDLE)
# actually place label to the left of chips row start; recompute start to leave room
start = (13.333 - total) / 2
for i, name in enumerate(chips):
    cx = start + i * (chip_w + chip_gap)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(cx), Inches(chip_y), Inches(chip_w), Inches(chip_h))
    solid(pill, PANEL)
    line(pill, TEAL, 1.25)
    pill.adjustments[0] = 0.5
    ptf = pill.text_frame
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ptf.margin_top = 0; ptf.margin_bottom = 0
    pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    pr = pp.add_run(); pr.text = name
    pr.font.name = FONT; pr.font.size = Pt(13); pr.font.bold = True; pr.font.color.rgb = WHITE

# hide the stray label (we won't use it) -> instead remove by moving offscreen
label.left = Inches(-5)

# ---- logo (top-right) ----
import os
logo_path = r"c:\Users\m.rahman\src\induction\logo-wimmera-white.png"
if os.path.exists(logo_path):
    logo_h = 1.2
    logo_w_guess = logo_h * (1024 / 750)
    lx = 13.333 - 0.6 - logo_w_guess
    slide.shapes.add_picture(logo_path, Inches(lx), Inches(0.42), height=Inches(logo_h))

# ---- footer bar ----
fh = 0.92
fy = 7.5 - fh
footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(fy), prs.slide_width, Inches(fh))
solid(footer, FOOTER_BG)
no_line(footer)
faccent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(fy), prs.slide_width, Inches(0.05))
solid(faccent, TEAL)
no_line(faccent)

ftxt = slide.shapes.add_textbox(Inches(0.65), Inches(fy), Inches(7.0), Inches(fh))
ftf = ftxt.text_frame
ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
ftf.word_wrap = True
fp = ftf.paragraphs[0]
fp.alignment = PP_ALIGN.LEFT
parts = [
    ("ICT108", TEAL, True),
    ("    ·    Guest Lecture    ·    ", MUTED, False),
    ("Tuesday 14 July 2026, 2–3 PM", WHITE, False),
]
for text, color, bold in parts:
    r = fp.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = color; r.font.bold = bold

frtxt = slide.shapes.add_textbox(Inches(7.0), Inches(fy), Inches(5.68), Inches(fh))
frtf = frtxt.text_frame
frtf.vertical_anchor = MSO_ANCHOR.MIDDLE
frtf.word_wrap = True
frp = frtf.paragraphs[0]
frp.alignment = PP_ALIGN.RIGHT
r = frp.add_run(); r.text = "Dr Mohammad Arifur Rahman"
r.font.name = FONT; r.font.size = Pt(13.5); r.font.color.rgb = TEAL; r.font.bold = True
frp2 = frtf.add_paragraph()
frp2.alignment = PP_ALIGN.RIGHT
r = frp2.add_run(); r.text = "Analyst Programmer  ·  Wimmera CMA, Victoria"
r.font.name = FONT; r.font.size = Pt(10.5); r.font.color.rgb = MUTED; r.font.bold = False

out = r"c:\Users\m.rahman\src\induction\lecture\flyer.pptx"
import os
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved", out)
